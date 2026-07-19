import io
import os
import datetime
from datetime import timedelta
from firebase_functions import https_fn, storage_fn, scheduler_fn
from firebase_admin import initialize_app, firestore, storage
import google.cloud.bigquery as bigquery
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import pandas as pd
import requests

# Initialize Firebase Admin SDK
firebase_app = initialize_app()
db = firestore.client()

def get_company_plan_limits(user_email):
    """Retrieve company plan limits (max_daily, max_kb) based on whitelist_users email."""
    try:
        user_ref = db.collection('whitelist_users').document(user_email.lower())
        user_snap = user_ref.get()
        if user_snap.exists:
            user_data = user_snap.to_dict()
            plan_name = user_data.get('plan', 'starter')
            
            plan_snap = db.collection('plan').document(plan_name).get()
            if plan_snap.exists:
                plan_data = plan_snap.to_dict()
                return {
                    'max_daily': plan_data.get('max_daily_photos', 300 if plan_name == 'pro' else 100),
                    'max_kb': plan_data.get('max_file_size_kb', 1024 if plan_name == 'pro' else 300)
                }
            return {
                'max_daily': 300 if plan_name == 'pro' else 100,
                'max_kb': 1024 if plan_name == 'pro' else 300
            }
    except Exception as e:
        print(f"Error checking plan limits: {e}")
    return {'max_daily': 100, 'max_kb': 300}  # Starter fallback

@https_fn.on_call()
def validateUpload(req: https_fn.CallableRequest) -> dict:
    """HTTPS Callable for client-side fast verification of files before upload."""
    if not req.auth:
        return {"error": "UNAUTHENTICATED"}
    
    data = req.data
    company_id = data.get('company_id')
    size_kb = data.get('size_kb', 0)
    
    if not company_id:
        return {"error": "MISSING_COMPANY_ID"}

    user_email = req.auth.token.get('email', '')
    limits = get_company_plan_limits(user_email)
    
    # 1. Size check
    if size_kb > limits['max_kb']:
        return {"error": "FILE_TOO_LARGE"}

    # 2. Daily count check
    today_str = datetime.date.today().isoformat()
    photos_ref = db.collection('photos')
    query_photos = photos_ref.where('company_id', '==', company_id) \
                             .where('upload_date', '==', today_str) \
                             .where('status', '==', 'done')
    
    # Perform count
    count = len(query_photos.get())
    if count >= limits['max_daily']:
        return {"error": "DAILY_LIMIT_REACHED"}
        
    return {"valid": True}

@https_fn.on_call()
def getSignedUploadUrl(req: https_fn.CallableRequest) -> dict:
    """
    Returns a short-lived GCS Signed URL for direct browser PUT upload.
    The client browser PUTs the image binary directly to GCS, bypassing
    Firebase Functions middleware entirely for maximum throughput.

    Input:  { gcs_path: str, content_type: str }
    Output: { signed_url: str }
    """
    if not req.auth:
        return {"error": "UNAUTHENTICATED"}

    data = req.data
    gcs_path = data.get('gcs_path', '')
    content_type = data.get('content_type', 'image/jpeg')

    if not gcs_path or not gcs_path.startswith('projects/'):
        return {"error": "INVALID_GCS_PATH"}

    try:
        bucket = storage.bucket()
        blob = bucket.blob(gcs_path)

        # Generate a PUT signed URL valid for 15 minutes
        signed_url = blob.generate_signed_url(
            expiration=datetime.timedelta(minutes=15),
            method='PUT',
            content_type=content_type,
            version='v4'
        )
        return {"signed_url": signed_url}
    except Exception as e:
        print(f"Error generating signed URL: {e}")
        return {"error": "SIGNED_URL_GENERATION_FAILED", "detail": str(e)}

@storage_fn.on_object_finalized()
def onPhotoUpload(event: storage_fn.StorageEvent) -> None:
    """Authoritative trigger checking uploads inside GCS projects/ paths."""
    path = event.data.name
    if not path.startswith('projects/'):
        return

    # Parse projects/{project_id}/{date}/{filename}
    parts = path.split('/')
    if len(parts) < 4:
        return
        
    project_id = parts[1]
    date_str = parts[2]
    filename = parts[3]

    # Clean extension & name
    extension = filename.split('.')[-1].lower() if '.' in filename else ''
    allowed_types = ['jpeg', 'jpg', 'png', 'gif']
    mime_type = event.data.content_type or ''

    # Get corresponding photo doc from Firestore
    photo_id = f"{project_id}_{filename.replace('.', '_')}"
    photo_ref = db.collection('photos').document(photo_id)
    photo_snap = photo_ref.get()

    photo_data = {}
    if photo_snap.exists:
        photo_data = photo_snap.to_dict()
    
    company_id = photo_data.get('company_id', 'default_company')
    uploaded_by = photo_data.get('uploaded_by', '')

    # Fetch limits
    limits = get_company_plan_limits(uploaded_by)

    rejection_reason = None

    # Check 1: Extension/MIME check
    if extension not in allowed_types or not mime_type.startswith('image/'):
        rejection_reason = "UNSUPPORTED_FILE_TYPE"
    
    # Check 2: Size check (event.data.size is in bytes)
    size_kb = event.data.size / 1024
    if not rejection_reason and size_kb > limits['max_kb']:
        rejection_reason = "FILE_TOO_LARGE"

    # Check 3: Daily count check
    if not rejection_reason:
        photos_ref = db.collection('photos')
        query_photos = photos_ref.where('company_id', '==', company_id) \
                                 .where('upload_date', '==', date_str) \
                                 .where('status', '==', 'done')
        
        count = len(query_photos.get())
        if count >= limits['max_daily']:
            rejection_reason = "DAILY_LIMIT_REACHED"

    # Handle rejection vs success
    bucket = storage.bucket(event.data.bucket)
    if rejection_reason:
        # Delete invalid file from Storage
        try:
            blob = bucket.blob(path)
            blob.delete()
        except Exception as e:
            print(f"Error deleting invalid storage object: {e}")

        # Update or create Firestore doc as rejected
        photo_ref.set({
            'id': photo_id,
            'project_id': project_id,
            'company_id': company_id,
            'filename': filename,
            'status': 'rejected',
            'reason': rejection_reason,
            'upload_date': date_str,
            'uploaded_by': uploaded_by,
            'updated_at': datetime.datetime.utcnow().isoformat()
        }, merge=True)
    else:
        # Update Firestore doc as done
        photo_ref.set({
            'status': 'done',
            'size_kb': size_kb,
            'updated_at': datetime.datetime.utcnow().isoformat()
        }, merge=True)

@https_fn.on_call()
def exportPPTX(req: https_fn.CallableRequest) -> dict:
    """Generates PowerPoint presentation with photos and captions. Returns 24-hr GCS Signed URL."""
    if not req.auth:
        return {"error": "UNAUTHENTICATED"}

    data = req.data or {}
    project_id = data.get('project_id')
    if not project_id:
        return {"error": "MISSING_PROJECT_ID"}

    view_mode = data.get('view_mode') or 'Desktop'
    is_mobile_mode = view_mode == 'Mobile'

    # Check if photos_data was compiled and sent directly from frontend
    photos_data = data.get('photos_data') or data.get('photos')
    photos_list = []
    if photos_data and isinstance(photos_data, list) and len(photos_data) > 0:
        photos_list = photos_data
    else:
        photos_ref = db.collection('photos')
        query_photos = photos_ref.where('project_id', '==', project_id).where('status', '==', 'done').get()
        now_utc = datetime.datetime.now(datetime.timezone.utc)
        photos_list = []
        for photo_doc in query_photos:
            d = photo_doc.to_dict()
            exp = d.get('expires_at')
            if exp:
                try:
                    exp_dt = datetime.datetime.fromisoformat(exp.replace('Z', '+00:00'))
                    if exp_dt < now_utc:
                        continue
                except Exception:
                    pass
            photos_list.append(d)
    
    # Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 layout
    
    bucket = storage.bucket()

    for p_data in photos_list:
        gcs_path = p_data.get('gcs_path')
        caption = p_data.get('caption', '') or p_data.get('comments', '')

        # Download photo to memory buffer
        try:
            img_stream = None
            raw_b64 = p_data.get('base64') if is_mobile_mode else (p_data.get('annotatedBase64') or p_data.get('base64'))
            if raw_b64 and isinstance(raw_b64, str):
                b64_part = raw_b64.split(',', 1)[1] if ',' in raw_b64 else raw_b64
                try:
                    img_stream = io.BytesIO(base64.b64decode(b64_part))
                except Exception:
                    pass

            if not img_stream and gcs_path:
                blob = bucket.blob(gcs_path)
                img_bytes = blob.download_as_bytes()
                img_stream = io.BytesIO(img_bytes)

            if not img_stream:
                continue

            # Open image with Pillow to determine original aspect ratio
            img = Image.open(img_stream)
            img_w, img_h = img.size
            img_aspect = img_w / img_h

            # Max box limits (90% slide size or centered if Mobile mode)
            max_w = Inches(9.0)
            max_h = Inches(4.5) if not is_mobile_mode else Inches(5.2)
            max_aspect = max_w / max_h

            if img_aspect > max_aspect:
                fit_w = max_w
                fit_h = max_w / img_aspect
            else:
                fit_h = max_h
                fit_w = max_h * img_aspect

            left = (prs.slide_width - fit_w) / 2
            top = (prs.slide_height - fit_h - (Inches(0.5) if not is_mobile_mode else Inches(0))) / 2
            if top < Inches(0.2):
                top = Inches(0.2)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, left, top, width=fit_w, height=fit_h)

            if not is_mobile_mode:
                txBox = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.6), Inches(9.0), Inches(0.4))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = caption if caption else "No Caption"
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.alignment = 1 # Center align
        except Exception as e:
            print(f"Skipping image {gcs_path} due to error: {e}")

    # Save PowerPoint to GCS Exports folder
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)

    export_path = f"exports/{project_id}_report.pptx"
    export_blob = bucket.blob(export_path)
    export_blob.upload_from_file(pptx_io, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

    # Generate 24-hr GCS Signed URL
    signed_url = export_blob.generate_signed_url(
        expiration=datetime.timedelta(hours=24),
        method='GET'
    )

    return {"downloadUrl": signed_url}

def _generate_pdf_report(req: https_fn.CallableRequest) -> dict:
    """Queries BigQuery/Firestore table or uses frontend photos_data, and creates professional PDF report of photos."""
    data = req.data or {}
    project_id = data.get('project_id')
    if not project_id:
        return {"error": "MISSING_PROJECT_ID"}

    view_mode = data.get('view_mode') or 'Desktop'
    is_mobile_mode = view_mode == 'Mobile'

    rows = []
    
    # 0. Check if photos_data was compiled and sent directly from frontend
    photos_data = data.get('photos_data') or data.get('photos')
    if photos_data and isinstance(photos_data, list) and len(photos_data) > 0:
        for p in photos_data:
            recs = p.get('recommendations_json') or p.get('recommendations') or p.get('recommendation') or []
            rec_text = '\n'.join(recs) if isinstance(recs, list) else str(recs)
            if not rec_text or rec_text == '[]':
                rec_text = "No specific recommendation noted."
            risk = p.get('risk') or p.get('risk_level') or "COMPLIANT"
            if "[CRITICAL]" in rec_text: risk = "CRITICAL"
            elif "[MAJOR]" in rec_text: risk = "MAJOR"
            elif "[MINOR]" in rec_text: risk = "MINOR"

            rows.append({
                'Photo Filename': p.get('filename') or 'IMG.jpg',
                'Caption': p.get('caption') or p.get('comments') or p.get('comments_text') or 'No visual defects observed.',
                'Recommendation': rec_text,
                'Risk Level': risk,
                'annotatedBase64': p.get('annotatedBase64') or p.get('base64') or '',
                'base64': p.get('base64') or '',
                'url': p.get('url') or p.get('localUrl') or p.get('thumbnailUrl') or '',
                'filename': p.get('filename') or 'IMG.jpg',
                'view_mode': view_mode
            })
    else:
        # 1. Try querying BigQuery first (as requested)
        try:
            bq_client = bigquery.Client()
            query_str = """
                SELECT 
                  JSON_VALUE(data, '$.filename') as filename,
                  JSON_VALUE(data, '$.caption') as caption,
                  JSON_VALUE(data, '$.annotatedBase64') as annotatedBase64,
                  JSON_VALUE(data, '$.base64') as base64,
                  JSON_VALUE(data, '$.url') as url,
                  JSON_VALUE(data, '$.recommendation') as recommendation,
                  JSON_VALUE(data, '$.risk') as risk
                FROM `mediaflow.photos_raw_latest` 
                WHERE JSON_VALUE(data, '$.project_id') = @project_id
            """
            job_config = bigquery.QueryJobConfiguration(
                query_parameters=[
                    bigquery.ScalarQueryParameter("project_id", "STRING", project_id)
                ]
            )
            query_job = bq_client.query(query_str, job_config=job_config)
            results = query_job.result()
            
            for r in results:
                rows.append({
                    'Photo Filename': r.filename or 'IMG.jpg',
                    'Caption': r.caption or 'No visual defects observed.',
                    'Recommendation': r.recommendation or 'No specific recommendation noted.',
                    'Risk Level': r.risk or 'COMPLIANT',
                    'annotatedBase64': r.annotatedBase64 or r.base64 or '',
                    'base64': r.base64 or '',
                    'url': r.url or '',
                    'filename': r.filename or 'IMG.jpg'
                })
        except Exception as bq_err:
            print(f"BigQuery query failed, falling back to Firestore: {bq_err}")
            photos_ref = db.collection('photos')
            query_photos = photos_ref.where('project_id', '==', project_id).where('status', '==', 'done').get()
            now_utc = datetime.datetime.now(datetime.timezone.utc)
            for doc_snap in query_photos:
                p_data = doc_snap.to_dict()
                exp = p_data.get('expires_at')
                if exp:
                    try:
                        exp_dt = datetime.datetime.fromisoformat(exp.replace('Z', '+00:00'))
                        if exp_dt < now_utc:
                            continue
                    except Exception:
                        pass
                recs = p_data.get('recommendations_json') or p_data.get('recommendations') or p_data.get('recommendation') or []
                rec_text = '\n'.join(recs) if isinstance(recs, list) else str(recs)
                if not rec_text or rec_text == '[]':
                    rec_text = "No specific recommendation noted."
                risk = p_data.get('risk') or p_data.get('risk_level') or "COMPLIANT"
                rows.append({
                    'Photo Filename': p_data.get('filename') or 'IMG.jpg',
                    'Caption': p_data.get('caption') or p_data.get('comments') or p_data.get('comments_text') or 'No visual defects observed.',
                    'Recommendation': rec_text,
                    'Risk Level': risk,
                    'annotatedBase64': p_data.get('annotatedBase64') or p_data.get('base64') or '',
                    'base64': p_data.get('base64') or '',
                    'url': p_data.get('url') or p_data.get('localUrl') or p_data.get('thumbnailUrl') or '',
                    'filename': p_data.get('filename') or 'IMG.jpg'
                })

    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    import base64
    import urllib.request
    from PIL import Image as PILImage

    def get_reportlab_image(photo_dict, max_width=100, max_height=75):
        """Fetch image from base64 buffer, local path, or URL, scale preserving aspect ratio, and return RLImage."""
        img_data = None
        b64_str = photo_dict.get('base64') if photo_dict.get('view_mode') == 'Mobile' else (photo_dict.get('annotatedBase64') or photo_dict.get('base64') or '')
        if b64_str and isinstance(b64_str, str):
            if ',' in b64_str:
                b64_str = b64_str.split(',', 1)[1]
            try:
                img_data = base64.b64decode(b64_str)
            except Exception as e:
                print(f"Failed to decode base64 for {photo_dict.get('Photo Filename')}: {e}")

        if not img_data:
            url = photo_dict.get('url') or photo_dict.get('localUrl') or photo_dict.get('thumbnailUrl') or ''
            if url and isinstance(url, str):
                if url.startswith('data:image/'):
                    try:
                        b64_part = url.split(',', 1)[1] if ',' in url else url
                        img_data = base64.b64decode(b64_part)
                    except Exception:
                        pass
                elif url.startswith('http://') or url.startswith('https://'):
                    try:
                        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
                        with urllib.request.urlopen(req, timeout=5) as resp:
                            img_data = resp.read()
                    except Exception as e:
                        print(f"Failed to fetch image from URL {url}: {e}")
                else:
                    if os.path.exists(url):
                        try:
                            with open(url, 'rb') as f:
                                img_data = f.read()
                        except Exception:
                            pass

        if not img_data:
            fname = photo_dict.get('Photo Filename') or photo_dict.get('filename') or ''
            if fname:
                for test_path in [fname, os.path.join('uploads', fname), os.path.join('/tmp', fname)]:
                    if os.path.exists(test_path):
                        try:
                            with open(test_path, 'rb') as f:
                                img_data = f.read()
                            break
                        except Exception:
                            pass

        if img_data:
            try:
                pil_img = PILImage.open(io.BytesIO(img_data))
                w, h = pil_img.size
                if w > 0 and h > 0:
                    ratio = min(max_width / float(w), max_height / float(h))
                    new_w = w * ratio
                    new_h = h * ratio
                    return RLImage(io.BytesIO(img_data), width=new_w, height=new_h)
            except Exception as e:
                print(f"Failed to process image dimensions with PIL/ReportLab: {e}")
                try:
                    return RLImage(io.BytesIO(img_data), width=max_width, height=max_height)
                except Exception:
                    pass
        return None

    pdf_io = io.BytesIO()
    doc = SimpleDocTemplate(pdf_io, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    normal_style = styles['Normal']
    normal_style.fontSize = 8
    normal_style.leading = 11

    title_style = ParagraphStyle(
        'ReportTitle',
        parent=styles['Heading1'],
        fontSize=16,
        leading=20,
        textColor=colors.HexColor('#111827')
    )

    story = [
        Paragraph("PT Safety Indonesia Utama - Assessment Report", title_style),
        Spacer(1, 4),
        Paragraph(f"Project: {project_name or project_id}", ParagraphStyle('Sub', parent=normal_style, fontSize=10, leading=13, textColor=colors.HexColor('#4B5563'), fontName='Helvetica-Bold')),
        Spacer(1, 14)
    ]

    header_style = ParagraphStyle(
        'TableHeader',
        parent=normal_style,
        fontSize=8,
        leading=10,
        textColor=colors.white,
        fontName='Helvetica-Bold'
    )

    table_data = []
    if is_mobile_mode:
        table_data.append([
            Paragraph("No.", header_style),
            Paragraph("Photo & Caption", header_style)
        ])
    else:
        table_data.append([
            Paragraph("No.", header_style),
            Paragraph("Photo Filename", header_style),
            Paragraph("Observation / Comments", header_style),
            Paragraph("Recommendation", header_style),
            Paragraph("Risk Level", header_style)
        ])

    for idx, r in enumerate(rows, 1):
        r['view_mode'] = view_mode
        rl_img = get_reportlab_image(r, max_width=450 if is_mobile_mode else 100, max_height=350 if is_mobile_mode else 75)
        fname_para = Paragraph(str(r.get('Photo Filename') or 'IMG.jpg'), normal_style)
        
        if is_mobile_mode:
            mobile_caption_style = ParagraphStyle(
                f'MobileCap_{idx}',
                parent=normal_style,
                alignment=1
            )
            caption_text = str(r.get('Caption') or r.get('Comments') or r.get('Photo Filename') or '')
            caption_para = Paragraph(caption_text, mobile_caption_style)
            if rl_img:
                photo_cell = [rl_img, Spacer(1, 8), caption_para]
            else:
                photo_cell = caption_para
            table_data.append([
                Paragraph(str(idx), normal_style),
                photo_cell
            ])
        else:
            if rl_img:
                photo_cell = [rl_img, Spacer(1, 4), fname_para]
            else:
                photo_cell = fname_para
            risk_val = str(r.get('Risk Level') or 'COMPLIANT')
            risk_style = ParagraphStyle(
                f'Risk_{idx}',
                parent=normal_style,
                fontName='Helvetica-Bold',
                alignment=1
            )
            if risk_val == 'CRITICAL':
                risk_style.textColor = colors.HexColor('#EF4444')
            elif risk_val == 'MAJOR':
                risk_style.textColor = colors.HexColor('#F97316')
            elif risk_val == 'MINOR':
                risk_style.textColor = colors.HexColor('#EAB308')
            else:
                risk_style.textColor = colors.HexColor('#10B981')

            table_data.append([
                Paragraph(str(idx), normal_style),
                photo_cell,
                Paragraph(str(r.get('Caption') or 'No visual defects observed.'), normal_style),
                Paragraph(str(r.get('Recommendation') or 'No specific recommendation noted.'), normal_style),
                Paragraph(risk_val, risk_style)
            ])

    if len(table_data) == 1:
        if is_mobile_mode:
            table_data.append([
                Paragraph("-", normal_style),
                Paragraph("No data found - No completed photos recorded.", normal_style)
            ])
        else:
            table_data.append([
                Paragraph("-", normal_style),
                Paragraph("No data found", normal_style),
                Paragraph("No completed photos recorded for this project.", normal_style),
                Paragraph("-", normal_style),
                Paragraph("-", normal_style)
            ])

    t = Table(table_data, colWidths=[40, 500] if is_mobile_mode else [25, 115, 170, 170, 60])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1F2937')),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E5E7EB')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F9FAFB')])
    ]))
    story.append(t)
    doc.build(story)
    pdf_io.seek(0)

    export_filename = data.get('export_filename') or f"{project_id}_report.pdf"
    bucket = storage.bucket()
    export_path = f"exports/{export_filename}"
    export_blob = bucket.blob(export_path)
    export_blob.upload_from_file(pdf_io, content_type='application/pdf')

    # Generate 24-hr GCS Signed URL
    signed_url = export_blob.generate_signed_url(
        expiration=datetime.timedelta(hours=24),
        method='GET',
        response_disposition=f'attachment; filename="{export_filename}"'
    )

    return {"downloadUrl": signed_url}

@https_fn.on_call()
def exportXLSX(req: https_fn.CallableRequest) -> dict:
    """Transitioned endpoint: previously mapped to xlsx, now outputs a PDF file as requested."""
    return _generate_pdf_report(req)

@https_fn.on_call()
def exportPDF(req: https_fn.CallableRequest) -> dict:
    """New endpoint: outputs a PDF file of the project report."""
    return _generate_pdf_report(req)

@scheduler_fn.on_schedule(schedule="every 24 hours")
def cleanup_exports(event: scheduler_fn.ScheduledEvent) -> None:
    """Cloud Scheduler job deleting export blobs older than 24 hours."""
    bucket = storage.bucket()
    blobs = bucket.list_blobs(prefix="exports/")
    
    now = datetime.datetime.now(datetime.timezone.utc)
    expiration_limit = timedelta(hours=24)

    for blob in blobs:
        # Ignore prefix folder itself
        if blob.name == "exports/":
            continue
            
        time_created = blob.time_created
        if now - time_created > expiration_limit:
            try:
                blob.delete()
                print(f"Deleted expired export blob: {blob.name}")
            except Exception as e:
                print(f"Failed to delete expired export blob {blob.name}: {e}")
